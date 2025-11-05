"""
文档处理服务
提供Word/PPT/图片等文档的文本提取和结构化处理
"""

import logging
import re
from typing import Dict, List, Any, Optional
from pathlib import Path
import io

logger = logging.getLogger(__name__)

# 尝试导入文档处理库（可选依赖）
try:
    from docx import Document as DocxDocument

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    logger.warning("python-docx 未安装，Word文档处理功能不可用")

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx 未安装，PPT文档处理功能不可用")

try:
    from PIL import Image
    import pytesseract
    import os

    # Windows 下自动检测 Tesseract-OCR 路径
    if os.name == "nt":  # Windows
        possible_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            os.path.join(
                os.environ.get("ProgramFiles", ""), "Tesseract-OCR", "tesseract.exe"
            ),
            os.path.join(
                os.environ.get("ProgramFiles(x86)", ""),
                "Tesseract-OCR",
                "tesseract.exe",
            ),
        ]

        for path in possible_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                logger.info(f"找到 Tesseract-OCR: {path}")
                break
        else:
            logger.warning("未找到 Tesseract-OCR 安装路径，OCR功能可能不可用")

    HAS_OCR = True
except ImportError:
    HAS_OCR = False
    logger.warning("PIL/pytesseract 未安装，图片OCR功能不可用")


class DocumentProcessingService:
    """文档处理服务"""

    def __init__(self, upload_dir: str = "uploads/expert_knowledge"):
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)

        # 支持的文档类型
        self.supported_types = {
            "word": [".doc", ".docx"],
            "ppt": [".ppt", ".pptx"],
            "image": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff"],
            "pdf": [".pdf"],
            "text": [".txt", ".md"],
        }

    async def extract_text_from_word(self, file_path: str) -> Dict[str, Any]:
        """从Word文档提取文本"""
        if not HAS_DOCX:
            raise ImportError("python-docx 未安装，无法处理Word文档")

        try:
            doc = DocxDocument(file_path)

            # 提取文本
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(
                        {
                            "text": para.text,
                            "style": para.style.name if para.style else None,
                            "is_heading": (
                                para.style.name.startswith("Heading")
                                if para.style
                                else False
                            ),
                        }
                    )

            # 提取表格
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)

            full_text = "\n".join([p["text"] for p in paragraphs])

            # 解析文档结构
            structure = self._parse_document_structure(paragraphs)

            return {
                "success": True,
                "full_text": full_text,
                "paragraphs": paragraphs,
                "tables": tables,
                "structure": structure,
                "word_count": len(full_text.split()),
            }

        except Exception as e:
            logger.error(f"提取Word文档文本失败: {e}")
            raise

    async def extract_text_from_ppt(self, file_path: str) -> Dict[str, Any]:
        """从PPT提取文本和图片说明"""
        if not HAS_PPTX:
            raise ImportError("python-pptx 未安装，无法处理PPT文档")

        try:
            prs = Presentation(file_path)

            slides = []
            full_text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_images = []

                # 提取文本（标题和内容）
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_text.append(text)
                        full_text_parts.append(text)

                        # 判断是否为标题（通常第一个文本框）
                        is_title = (
                            shape.placeholder_format.idx == 0
                            if hasattr(shape, "placeholder_format")
                            else False
                        )

                    # 提取图片说明（如果有）
                    if hasattr(shape, "image"):
                        image_info = {
                            "image_path": None,  # 实际应该提取图片并保存
                            "description": (
                                shape.name if hasattr(shape, "name") else None
                            ),
                        }
                        slide_images.append(image_info)

                slides.append(
                    {
                        "slide_number": slide_num,
                        "title": slide_text[0] if slide_text else None,
                        "content": (
                            "\n".join(slide_text[1:]) if len(slide_text) > 1 else None
                        ),
                        "full_text": "\n".join(slide_text),
                        "images": slide_images,
                    }
                )

            full_text = "\n\n".join(full_text_parts)

            return {
                "success": True,
                "full_text": full_text,
                "slides": slides,
                "slide_count": len(slides),
                "word_count": len(full_text.split()),
            }

        except Exception as e:
            logger.error(f"提取PPT文档文本失败: {e}")
            raise

    async def extract_text_from_image(
        self, file_path: str, language: str = "chi_sim+eng"
    ) -> Dict[str, Any]:
        """OCR识别图片中的文本"""
        if not HAS_OCR:
            raise ImportError("PIL/pytesseract 未安装，无法进行OCR识别")

        try:
            # 打开图片
            image = Image.open(file_path)

            # OCR识别
            text = pytesseract.image_to_string(image, lang=language)

            # 获取详细信息（包括位置信息）
            data = pytesseract.image_to_data(
                image, lang=language, output_type=pytesseract.Output.DICT
            )

            # 提取关键信息
            words = []
            for i in range(len(data["text"])):
                if data["text"][i].strip():
                    words.append(
                        {
                            "text": data["text"][i],
                            "confidence": (
                                int(data["conf"][i])
                                if data["conf"][i] != "-1"
                                else None
                            ),
                            "left": data["left"][i],
                            "top": data["top"][i],
                            "width": data["width"][i],
                            "height": data["height"][i],
                        }
                    )

            # 计算平均置信度
            confidences = [
                w["confidence"] for w in words if w["confidence"] is not None
            ]
            avg_confidence = (
                sum(confidences) / len(confidences) if confidences else None
            )

            return {
                "success": True,
                "full_text": text,
                "words": words,
                "word_count": len([w for w in words if w["text"].strip()]),
                "average_confidence": avg_confidence,
            }

        except Exception as e:
            logger.error(f"OCR识别图片文本失败: {e}")
            raise

    async def extract_text_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """从PDF提取文本（需要pdfplumber或PyPDF2）"""
        try:
            # 这里应该使用 pdfplumber 或 PyPDF2
            # 暂时返回占位符
            logger.warning("PDF处理功能待实现")

            return {
                "success": False,
                "message": "PDF处理功能待实现，需要安装 pdfplumber 或 PyPDF2",
            }

        except Exception as e:
            logger.error(f"提取PDF文本失败: {e}")
            raise

    async def extract_text_from_file(
        self, file_path: str, file_type: Optional[str] = None
    ) -> Dict[str, Any]:
        """根据文件类型自动提取文本"""
        try:
            file_path_obj = Path(file_path)

            # 自动检测文件类型
            if not file_type:
                file_type = self._detect_file_type(file_path)

            if file_type == "word":
                return await self.extract_text_from_word(file_path)
            elif file_type == "ppt":
                return await self.extract_text_from_ppt(file_path)
            elif file_type == "image":
                return await self.extract_text_from_image(file_path)
            elif file_type == "pdf":
                return await self.extract_text_from_pdf(file_path)
            elif file_type == "text":
                return await self._extract_text_from_text_file(file_path)
            else:
                raise ValueError(f"不支持的文件类型: {file_type}")

        except Exception as e:
            logger.error(f"提取文件文本失败: {e}")
            raise

    def parse_document_structure(
        self, paragraphs: List[Dict[str, Any]]
    ) -> Dict[str, Any]:
        """解析文档结构（标题、段落、列表）"""
        structure = {
            "title": None,
            "headings": [],
            "paragraphs": [],
            "lists": [],
            "sections": [],
        }

        current_section = None

        for para in paragraphs:
            text = para.get("text", "").strip()
            if not text:
                continue

            is_heading = para.get("is_heading", False)
            style = para.get("style", "")

            if is_heading or (style and "Heading" in style):
                # 这是一个标题
                heading_level = self._extract_heading_level(style)

                heading = {
                    "text": text,
                    "level": heading_level,
                    "position": len(structure["paragraphs"]),
                }
                structure["headings"].append(heading)

                # 创建新章节
                if current_section:
                    structure["sections"].append(current_section)

                current_section = {
                    "title": text,
                    "level": heading_level,
                    "paragraphs": [],
                }
            else:
                # 普通段落
                para_info = {"text": text}
                structure["paragraphs"].append(para_info)

                if current_section:
                    current_section["paragraphs"].append(para_info)

        # 添加最后一个章节
        if current_section:
            structure["sections"].append(current_section)

        # 设置文档标题（第一个标题）
        if structure["headings"]:
            structure["title"] = structure["headings"][0]["text"]

        return structure

    def extract_key_concepts(self, text: str, max_concepts: int = 10) -> List[str]:
        """提取关键概念（简单实现，基于词频）"""
        try:
            # 简单的关键词提取
            # 移除常见停用词
            stop_words = {
                "的",
                "是",
                "在",
                "了",
                "和",
                "有",
                "也",
                "都",
                "为",
                "与",
                "the",
                "is",
                "a",
                "an",
                "and",
                "or",
                "but",
                "in",
                "on",
                "at",
                "to",
                "for",
                "of",
                "with",
            }

            # 分词（简单空格分词，中文需要更复杂的处理）
            words = re.findall(r"\b\w+\b", text.lower())

            # 过滤停用词和短词
            filtered_words = [w for w in words if w not in stop_words and len(w) > 2]

            # 统计词频
            from collections import Counter

            word_freq = Counter(filtered_words)

            # 返回最常见的词
            concepts = [word for word, count in word_freq.most_common(max_concepts)]

            return concepts

        except Exception as e:
            logger.error(f"提取关键概念失败: {e}")
            return []

    def generate_summary(self, text: str, max_sentences: int = 3) -> str:
        """生成摘要（简单实现，取前几个句子）"""
        try:
            # 简单实现：按句号分割，取前几个句子
            sentences = re.split(r"[。！？.!?]", text)
            sentences = [s.strip() for s in sentences if s.strip()]

            summary_sentences = sentences[:max_sentences]
            summary = "。".join(summary_sentences)

            if len(sentences) > max_sentences:
                summary += "..."

            return summary

        except Exception as e:
            logger.error(f"生成摘要失败: {e}")
            return text[:200] + "..." if len(text) > 200 else text

    # 私有辅助方法

    def _detect_file_type(self, file_path: str) -> str:
        """检测文件类型"""
        file_path_obj = Path(file_path)
        extension = file_path_obj.suffix.lower()

        for file_type, extensions in self.supported_types.items():
            if extension in extensions:
                return file_type

        return "other"

    def _extract_text_from_text_file(self, file_path: str) -> Dict[str, Any]:
        """从文本文件提取文本"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()

            return {"success": True, "full_text": text, "word_count": len(text.split())}

        except Exception as e:
            logger.error(f"读取文本文件失败: {e}")
            raise

    def _parse_document_structure(
        self, paragraphs: List[Dict[str, Any]]
    ) -> Dict[str, Any]:
        """解析文档结构（内部方法）"""
        return self.parse_document_structure(paragraphs)

    def _extract_heading_level(self, style: str) -> int:
        """提取标题级别"""
        # 从样式名中提取级别（如 "Heading 1" -> 1）
        match = re.search(r"heading\s*(\d+)", style.lower())
        if match:
            return int(match.group(1))
        return 1  # 默认级别
